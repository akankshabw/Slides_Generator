import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from newspaper import Article
import cohere
import os
import re
from dotenv import load_dotenv
import argparse
from urllib.parse import urlparse, parse_qs, unquote
from pptx.oxml.xmlchemy import OxmlElement
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
# Load API key from .env file
load_dotenv()
co = cohere.Client(os.getenv("COHERE_API_KEY"))

# --------------------------
# Step 1: Web Search (DuckDuckGo)
# --------------------------
def search_duckduckgo(query, num_results=DEFAULT_NUM_RESULTS):
    url = f"https://duckduckgo.com/html/?q={query}"
    headers = {
        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/116.0.5845.141 Safari/537.36"
    }
    res = requests.get(url, headers=headers)
    soup = BeautifulSoup(res.text, "html.parser")
    raw_links = [a['href'] for a in soup.select('.result__a')][:num_results]

    cleaned_links = []
    for raw in raw_links:
        parsed = urlparse(raw)
        query = parse_qs(parsed.query)
        if 'uddg' in query:
            real_url = unquote(query['uddg'][0])
            cleaned_links.append(real_url)

    return cleaned_links

# --------------------------
# Step 2: Extract article text using newspaper3k
# --------------------------
def scrape_and_extract_text(url):
    try:
        article = Article(url)
        article.download()
        article.parse()
        text = article.text
        print(f"\n=== Extracted from {url} ===")
        print(text[:1000])
        print("===================================\n")
        return text[:MAX_ARTICLE_LENGTH]  
    except Exception as e:
        print(f"❌ Failed to extract from {url}: {e}")
        return ""

# --------------------------
# Step 3: Summarize via Cohere
# --------------------------
def summarize_content(text, topic):
    prompt = f"You are given the following text which is information gathered from internet. Use this and your own knowledge and create 20-25 bullet points for a presentation on '{topic}':\n\n{text}"
    response = co.chat(model="command-r-plus", message=prompt)
    return response.text

# --------------------------
# Step 4: Generate Slide Content
# --------------------------
def generate_slide_deck(topic, summarized_points):
    prompt = f"""
You are a professional presentation assistant. Use the following context to create a 7-slide presentation on "{topic}".

--- CONTEXT START ---
{summarized_points}
--- CONTEXT END ---

Slide Structure:
1. Title Slide (only the topic)
2. Overview
3–6: Key ideas / trends / arguments  (each slide = title + 6-8 bullet points)
7. Conclusion

Guidelines:
- Use concise, impactful titles.
- Do not include slide numbers or any extra text or ###.
- Output titles and bullet points clearly.
- The ttle slide should contain the topic in the middle of the slide.
- Give 6-8 bullet points for slides 3-6.
"""
    response = co.chat(model="command-r-plus", message=prompt)
    return response.text

# --------------------------
# Step 5: Create PowerPoint File
# --------------------------
def create_ppt_from_text(slide_text, filename="output.pptx"):
    prs = Presentation()
    slides = slide_text.strip().split("\n\n")

    for idx, slide in enumerate(slides):
        if idx >= 8:  # Limit to 7 slides
            break
        lines = [line.strip() for line in slide.strip().split("\n") if line.strip()]
        if not lines:
            continue

        raw_title = lines[0]
        title = re.sub(r'^[#>\-\*\s]+', '', raw_title).strip()
        title = re.sub(r'[\*\_`]+$', '', title).strip()  # remove trailing **, __, or ` if present
        title = re.sub(r'^Slide\s*\d+:\s*', '', title, flags=re.IGNORECASE)
        content = lines[1:]

        # Use Title Only layout for first slide
        slide_layout = prs.slide_layouts[5] if idx == 0 else prs.slide_layouts[1]
        slide_obj = prs.slides.add_slide(slide_layout)

        # Set background color (soft light blue)
        slide_obj.background.fill.solid()
        slide_obj.background.fill.fore_color.rgb = RGBColor(BACKGROUND_COLOR[0], BACKGROUND_COLOR[1], BACKGROUND_COLOR[2])

        # Title formatting
        title_shape = slide_obj.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(38)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 30, 80)
        title_frame = title_shape.text_frame

        if idx == 0:
            # Enlarge and center the title box
            title_shape.top = Inches(2)
            title_shape.left = Inches(1)
            title_shape.height = Inches(2)
            title_shape.width = Inches(8)
            title.font.name = TITLE_FONT

            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = title_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(44)
            para.font.bold = True
            para.font.color.rgb = RGBColor(30, 30, 80)

        # Content for non-title slides
        if idx != 0 and content:
            # Add a thin horizontal line under the title
            left = Inches(0.5)
            top = Inches(1.3)
            width = Inches(8.5)
            height = Inches(0.05)

            shape = slide_obj.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 60, 120)  # deep blue accent
            shape.line.fill.background()  # remove border

            content_box = slide_obj.placeholders[1]
            tf = content_box.text_frame
            tf.clear()

            for i, bullet in enumerate(content):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                clean_bullet = re.sub(r"^\s*-\s*", "• ", bullet)
                p.text = clean_bullet
                p.level = 0
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(40, 40, 40)
                p.font.name = TITLE_FONT
                p.space_after = Pt(6)

                # Remove PowerPoint's automatic bullets
                pPr = p._element.get_or_add_pPr()
                buNone = OxmlElement('a:buNone')
                pPr.insert(0, buNone)


    prs.save(filename)
    print(f"✅ Aesthetic presentation saved to {filename}")

# --------------------------
# Master Function
# --------------------------
def generate_slide_deck_for_topic(topic):
    print(f"\n Searching web for: {topic}")
    urls = search_duckduckgo(topic)
    print(" Found URLs:", urls)

    print("\n Extracting content from web...")
    all_text = ""
    for url in urls:
        extracted = scrape_and_extract_text(url)
        all_text += extracted + "\n"

    if not all_text.strip():
        print(" No usable text extracted from web. Exiting.")
        return

    print("\n Summarizing content...")
    summary = summarize_content(all_text, topic)
    print(" Summary:\n", summary)

    print("\n Generating slides...")
    slide_text = generate_slide_deck(topic, summary)
    print(" Slide Content:\n", slide_text)

    print("\n Creating PowerPoint...")
    create_ppt_from_text(slide_text)

# --------------------------
# Run from Command Line
# --------------------------
if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    user_topic = input("Enter the topic for the slide deck: ")
    generate_slide_deck_for_topic(user_topic)
