import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from slide_generator.config import BACKGROUND_COLOR 

def create_ppt_from_text(slide_text, filename="output.pptx"):
    prs = Presentation()
    slides = slide_text.strip().split("\n\n")

    for idx, slide in enumerate(slides):
        if idx >= 7:  # Limit to 7 slides
            break
        lines = [line.strip() for line in slide.strip().split("\n") if line.strip()]
        if not lines:
            continue

        raw_title = lines[0]
        title = re.sub(r'^[#>\-\*\s]+', '', raw_title).strip()
        title = re.sub(r'[\*\_`]+$', '', title).strip() 
        title = re.sub(r'^Slide\s*\d+:\s*', '', title, flags=re.IGNORECASE)
        title = re.sub(r'^Key Idea\s*\d+\s*:\s*', '', title, flags=re.IGNORECASE)
        content = lines[1:]

        # Use Title Only layout for first slide
        slide_layout = prs.slide_layouts[5] if idx == 0 else prs.slide_layouts[1]
        slide_obj = prs.slides.add_slide(slide_layout)

        slide_obj.background.fill.solid()
        slide_obj.background.fill.fore_color.rgb = RGBColor(BACKGROUND_COLOR[0], BACKGROUND_COLOR[1], BACKGROUND_COLOR[2])

        # Title formatting
        title_shape = slide_obj.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(38)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 30, 80)
        title_frame = title_shape.text_frame

        if idx == 0:
            # Enlarge and center the title box
            title_shape.top = Inches(2)
            title_shape.left = Inches(1)
            title_shape.height = Inches(2)
            title_shape.width = Inches(8)
   

            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = title_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(44)
            para.font.bold = True
            para.font.color.rgb = RGBColor(30, 30, 80)

        # Content for non-title slides
        if idx != 0 and content:
            left = Inches(0.5)
            top = Inches(1.6)
            width = Inches(9.0)
            height = Inches(0.05)

            shape = slide_obj.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 60, 120) 
            shape.line.fill.background()  

            content_box = slide_obj.placeholders[1]
            tf = content_box.text_frame
            tf.clear()

            for i, bullet in enumerate(content):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                clean_bullet = re.sub(r"^\s*-\s*", "• ", bullet)
                p.text = clean_bullet
                p.level = 0
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(40, 40, 40)
           
                p.space_after = Pt(6)

                pPr = p._element.get_or_add_pPr()
                buNone = OxmlElement('a:buNone')
                pPr.insert(0, buNone)


    prs.save(filename)
    print(f"Presentation saved as {filename}")
